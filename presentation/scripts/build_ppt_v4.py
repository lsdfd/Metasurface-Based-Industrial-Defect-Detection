from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "presentation" / "build"
FORMULA_DIR = ROOT / "presentation" / "generated_figures" / "formulas"
TABLE_DIR = ROOT / "presentation" / "generated_figures" / "tables"
AI_V2 = ROOT / "presentation" / "ai_images_v2"
AI_V3 = ROOT / "presentation" / "ai_images_v3"
AI_V4 = ROOT / "presentation" / "ai_images_v4"
FABRIC = ROOT / "fabric_defect_detection-main" / "paper_assets"
DAGM = ROOT / "mixed-segdec-net-comind2021-master" / "paper_assets"
PPTX_PATH = OUT_DIR / "metasurface_industrial_defect_detection_v4.pptx"

NAVY = RGBColor(22, 58, 108)
DARK = RGBColor(36, 50, 74)
MID = RGBColor(95, 111, 132)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(237, 243, 252)
LINE = RGBColor(213, 222, 236)


def ensure_dirs():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    FORMULA_DIR.mkdir(parents=True, exist_ok=True)
    TABLE_DIR.mkdir(parents=True, exist_ok=True)


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, title: str, section: str):
    add_text(slide, Inches(0.55), Inches(0.25), Inches(10.6), Inches(0.48), title, size=21, bold=True, color=NAVY)
    add_text(slide, Inches(11.0), Inches(0.29), Inches(1.75), Inches(0.3), section, size=9, color=MID, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.88), Inches(12.25), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, text: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.86), Inches(12.25), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE
    shape.line.color.rgb = LINE
    add_text(slide, Inches(0.72), Inches(6.89), Inches(11.9), Inches(0.31), text, size=12, bold=True, color=NAVY)


def blank(prs, title: str, footer: str, section: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, title, section)
    add_footer(slide, footer)
    return slide


def add_image_fit(slide, path: Path, x, y, w, h):
    from PIL import Image

    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        fw = w
        fh = w / img_ratio
    else:
        fh = h
        fw = h * img_ratio
    slide.shapes.add_picture(str(path), x + (w - fw) / 2, y + (h - fh) / 2, width=fw, height=fh)


def add_image_frame(slide, path: Path, x, y, w, h):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    add_image_fit(slide, path, x + Inches(0.06), y + Inches(0.06), w - Inches(0.12), h - Inches(0.12))


def formula(name: str, expr: str, *, fontsize=25, width=8.0, height=0.75) -> Path:
    path = FORMULA_DIR / f"{name}.png"
    fig = plt.figure(figsize=(width, height), dpi=220)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, expr, ha="center", va="center", fontsize=fontsize, color="#163A6C")
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    return path


def table_image(name: str, headers: list[str], rows: list[list[str]], *, figsize=(10, 2.2), fontsize=10) -> Path:
    path = TABLE_DIR / f"{name}.png"
    fig, ax = plt.subplots(figsize=figsize, dpi=220)
    ax.axis("off")
    table = ax.table(cellText=rows, colLabels=headers, loc="center", cellLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(fontsize)
    table.scale(1, 1.55)
    for (r, c), cell in table.get_celld().items():
        cell.set_edgecolor("#D5DEEC")
        if r == 0:
            cell.set_facecolor("#163A6C")
            cell.get_text().set_color("white")
            cell.get_text().set_weight("bold")
        else:
            cell.set_facecolor("#F8FAFD" if r % 2 else "white")
            cell.get_text().set_color("#24324A")
    fig.savefig(path, transparent=False, bbox_inches="tight", pad_inches=0.08, facecolor="white")
    plt.close(fig)
    return path


def image_slide(prs, title, image, footer, section, *, formula_path: Path | None = None):
    slide = blank(prs, title, footer, section)
    if formula_path:
        add_image_frame(slide, image, Inches(0.65), Inches(1.08), Inches(12.05), Inches(4.78))
        add_image_fit(slide, formula_path, Inches(1.5), Inches(5.95), Inches(10.35), Inches(0.55))
    else:
        add_image_frame(slide, image, Inches(0.65), Inches(1.08), Inches(12.05), Inches(5.55))


def two_image_slide(prs, title, left, right, footer, section):
    slide = blank(prs, title, footer, section)
    add_image_frame(slide, left, Inches(0.65), Inches(1.15), Inches(5.85), Inches(5.35))
    add_image_frame(slide, right, Inches(6.85), Inches(1.15), Inches(5.85), Inches(5.35))


def three_image_slide(prs, title, a, b, c, footer, section):
    slide = blank(prs, title, footer, section)
    add_image_frame(slide, a, Inches(0.55), Inches(1.35), Inches(4.05), Inches(4.75))
    add_image_frame(slide, b, Inches(4.65), Inches(1.35), Inches(4.05), Inches(4.75))
    add_image_frame(slide, c, Inches(8.75), Inches(1.35), Inches(4.05), Inches(4.75))


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_text(slide, Inches(1.0), Inches(1.55), Inches(11.25), Inches(0.85), "超表面光学前端辅助的工业缺陷检测方法研究", size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.35), Inches(2.55), Inches(10.6), Inches(0.4), "Metasurface-based Optical Frontend for Industrial Defect Detection", size=16, color=DARK, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.3), Inches(3.35), Inches(8.75), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    add_text(slide, Inches(4.5), Inches(4.15), Inches(4.3), Inches(0.45), "李升欣", size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4.2), Inches(4.78), Inches(4.9), Inches(0.4), "2026年5月", size=15, color=DARK, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(6.5), Inches(10.9), Inches(0.35), "Fabric 二分类验证  |  DAGM / SegDecNet 分割闭环  |  Kernel-to-PSF 物理映射", size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)


def build():
    ensure_dirs()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    kd_formula = formula("v3_kd_loss", r"$\mathcal{L}=\mathcal{L}_{task}+\lambda_{cls}\mathcal{L}_{KD}^{cls}+\lambda_{seg}\mathcal{L}_{KD}^{seg}+\lambda_{vol}\mathcal{L}_{KD}^{vol}$")
    split_formula = formula("v3_kernel_split", r"$K^{+}=\max(K,0),\quad K^{-}=\max(-K,0),\quad K=K^{+}-K^{-}$")
    fabric_table = table_image(
        "v3_fabric_metrics",
        ["Metric", "Teacher", "R1 student @0.5", "R1 student best"],
        [["F1", "≈0.975", "0.2353", "0.8571"], ["Precision / Recall", "-", "-", "0.75 / 1.00"], ["Threshold", "-", "0.50", "≈0.90"]],
        figsize=(9.2, 1.55),
        fontsize=9,
    )
    dagm_table = table_image(
        "v3_dagm_metrics",
        ["AP", "AUC", "IoU", "Dice", "Precision", "Recall"],
        [["1.000", "1.000", "0.9145", "0.9349", "0.9958", "0.9176"]],
        figsize=(9.4, 1.25),
        fontsize=10,
    )
    comparison_table = table_image(
        "v3_case_comparison",
        ["Case", "Defect scene", "Output", "Main result", "Role"],
        [["Fabric", "textile patch", "binary OK/NG", "F1=0.8571", "first demo"], ["DAGM", "texture anomaly", "score + mask", "IoU=0.9145", "main closed loop"]],
        figsize=(10.8, 1.55),
        fontsize=8.5,
    )

    cover(prs)
    image_slide(prs, "背景：工业检测计算瓶颈与超表面视觉前端机会", AI_V4 / "v4_background_infographic.png", "工业质检需要高速低功耗视觉前端，超表面提供了把早期卷积前移到成像链路的入口。", "背景")
    image_slide(prs, "总体方法：传统电子模型到光电融合 student", AI_V4 / "v4_methodology_overview.png", "主线是 teacher 蒸馏、optical student、kernel-to-PSF 和轻量电子后端。", "方法")
    image_slide(prs, "光电融合系统：metasurface + CMOS + electronic backend", AI_V4 / "v4_hybrid_optical_architecture.png", "超表面承担早期光学卷积，电子后端负责 calibration、非线性和任务输出。", "方法")
    image_slide(prs, "知识蒸馏机制：对齐 score、mask 与 feature volume", AI_V4 / "v4_distillation_mechanism.png", "蒸馏把强电子 teacher 的任务能力迁移到受物理约束的 optical student。", "方法", formula_path=kd_formula)
    image_slide(prs, "Kernel-to-PSF：从 signed kernel 到超表面可实现目标", AI_V4 / "v4_kernel_to_psf.png", "正负拆分是把含符号卷积核映射到非负光强 PSF 的关键步骤。", "方法", formula_path=split_formula)
    image_slide(prs, "两个案例：Fabric 二分类与 DAGM 像素级分割", AI_V4 / "v4_case_overview.png", "Fabric 是方法起点，DAGM 是高性能分割和物理映射闭环。", "方法")

    image_slide(prs, "Fabric / AITEX：织物 patch 缺陷二分类任务", AI_V4 / "v4_case_overview.png", "Fabric 检测织物 patch 是否存在缺陷，不输出具体缺陷位置。", "Fabric")
    image_slide(prs, "Fabric teacher/student 架构", AI_V4 / "v4_fabric_architecture.png", "R1 student 使用 64×64 输入和一层 optical convolution bank 学习二分类判别。", "Fabric")
    two_image_slide(prs, "Fabric 主结果与阈值校准", FABRIC / "figures_main" / "results" / "fabric_teacher_student_summary_bar.png", FABRIC / "figures_process" / "threshold_sweep" / "fabric_r1_threshold_story.png", "R1 student 在最佳阈值下 F1=0.8571，但部署前必须做阈值校准。", "Fabric")
    image_slide(prs, "Fabric 指标表：teacher、默认阈值和最佳阈值", fabric_table, "默认阈值会低估 student，最佳阈值下模型已具备作为 optical frontend 原型的价值。", "Fabric")
    three_image_slide(prs, "Fabric optical kernels：signed / positive / negative", FABRIC / "figures_main" / "kernels" / "kernel_grid_signed.png", FABRIC / "figures_main" / "kernels" / "kernel_grid_positive.png", FABRIC / "figures_main" / "kernels" / "kernel_grid_negative.png", "Fabric kernels 已完成正负拆分，可进入 metasurface PSF target 链路。", "Fabric")
    image_slide(prs, "Fabric 计算量节省", FABRIC / "figures_main" / "compute" / "fabric_compute_comparison_bar.png", "Fabric teacher 到 student 约 175× 参数减少、约 99.5× MAC 减少。", "Fabric")
    image_slide(prs, "Fabric 补充资产：teacher confusion matrix / demo UI", FABRIC / "figures_main" / "results" / "fabric_teacher_confusion_matrix.png", "Fabric 补充图用于说明原二分类链路和可视化界面，不作为最终主结果。", "Fabric")

    image_slide(prs, "DAGM Class7：工业纹理缺陷分类 + 像素级分割", AI_V4 / "v4_case_overview.png", "DAGM 不只判断有没有缺陷，还输出 pixel-level defect mask，是当前主结果。", "DAGM")
    image_slide(prs, "DAGM teacher/student 架构", AI_V4 / "v4_dagm_architecture.png", "Student 保留 SegDecNet-style 电子后端，把早期 volume 前端替换为 optical convolution bank。", "DAGM")
    image_slide(prs, "DAGM 两阶段蒸馏训练", AI_V4 / "v4_two_stage_training.png", "先强化 volume/segmentation 对齐，再联合优化分类、分割和任务损失。", "DAGM", formula_path=kd_formula)
    image_slide(prs, "DAGM 定量结果", DAGM / "figures_main" / "metrics" / "dagm_full_validation_bar.png", "Full validation 上 AP/AUC=1.0，IoU=0.9145，Dice=0.9349。", "DAGM")
    image_slide(prs, "DAGM 指标表", dagm_table, "Precision 高于 Recall，说明 mask 偏保守但误检很少。", "DAGM")
    image_slide(prs, "DAGM mask 定性结果", DAGM / "figures_main" / "qualitative_masks" / "mask_visualization_contact_sheet_12.jpg", "预测热图基本落在真实缺陷区域，说明 student 学到了有效定位能力。", "DAGM")
    image_slide(prs, "DAGM threshold sweep", DAGM / "figures_process" / "threshold_sweep" / "dagm_threshold_sweep.png", "DAGM 默认 threshold=0.5 已稳定，结果不是靠手动调阈值刷出来的。", "DAGM")
    three_image_slide(prs, "DAGM optical kernels：signed / positive / negative", DAGM / "figures_main" / "kernels" / "kernel_grid_signed.png", DAGM / "figures_main" / "kernels" / "kernel_grid_positive.png", DAGM / "figures_main" / "kernels" / "kernel_grid_negative.png", "DAGM kernels 呈现纹理、边缘和方向响应，是 PSF 设计的核心输入。", "DAGM")
    two_image_slide(prs, "DAGM PSF target 与 backphase", DAGM / "figures_main" / "psf_targets" / "psf_target_center_crop.png", DAGM / "figures_process" / "metasurface_probe" / "psf_backphase_preview.png", "Learned kernels 已经可转成 target PSF 和初始 backphase。", "DAGM")
    two_image_slide(prs, "Metasurface feasibility probe", DAGM / "figures_main" / "metasurface_probe" / "kernel00_positive_probe.png", DAGM / "figures_main" / "metasurface_probe" / "metasurface_probe_cosine_bar.png", "代表 kernel 的 PSF 拟合 cosine similarity 达到 0.979-0.991。", "DAGM")
    image_slide(prs, "Metasurface probe 补充样例", DAGM / "figures_main" / "metasurface_probe" / "kernel37_negative_probe.png", "不平衡 kernel 的 negative branch 也能作为压力测试样例进入可行性分析。", "DAGM")
    image_slide(prs, "DAGM 计算量节省", DAGM / "figures_main" / "compute" / "compute_comparison_bar.png", "Hybrid electronic backend 相对 teacher@256 理论电子 MACs 减少约 905×。", "DAGM")
    image_slide(prs, "DAGM 计算量表", table_image("v3_dagm_compute", ["Comparison", "Value"], [["Teacher @256", "21.08G MACs"], ["Hybrid backend", "23.3M electronic MACs"], ["Reduction @256", "905× fewer"], ["Reduction @512", "3618× fewer"]], figsize=(7.6, 2.2), fontsize=10), "这些是理论电子 MAC reduction，实际速度还取决于曝光、读出、ADC 和硬件实现。", "DAGM")

    image_slide(prs, "项目闭环总结", AI_V4 / "v4_project_summary.png", "本课题已经形成 teacher -> optical student -> kernels -> PSF -> metasurface probe -> compute reduction 的证据链。", "总结")
    image_slide(prs, "两个案例对照总结", comparison_table, "Fabric 验证二分类方法起点，DAGM 验证高性能分割和物理映射闭环。", "总结")
    image_slide(prs, "未来展望：从分割到定位检测", AI_V4 / "v4_future_roadmap.png", "下一步可扩展到 YOLO、PCB、wafer、chip defect localization，并加入 hardware-aware retraining。", "展望")

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build()
