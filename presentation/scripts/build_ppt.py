from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "presentation" / "build"
FIG_DIR = ROOT / "presentation" / "generated_figures"
AI_DIR = ROOT / "presentation" / "ai_images"
FABRIC = ROOT / "fabric_defect_detection-main" / "paper_assets"
DAGM = ROOT / "mixed-segdec-net-comind2021-master" / "paper_assets"

PPTX_PATH = OUT_DIR / "metasurface_industrial_defect_detection_v1.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_Y = Inches(0.28)
TITLE_X = Inches(0.55)
TITLE_W = Inches(12.2)
LINE_Y = Inches(0.86)
CONTENT_TOP = Inches(1.08)
CONTENT_H = Inches(5.55)
FOOTER_Y = Inches(6.84)

NAVY = RGBColor(22, 58, 108)
BLUE = RGBColor(47, 107, 255)
TEAL = RGBColor(42, 157, 143)
RED = RGBColor(231, 111, 81)
GRAY = RGBColor(245, 247, 250)
DARK = RGBColor(36, 50, 74)
MID = RGBColor(100, 116, 139)
LIGHT_LINE = RGBColor(213, 222, 236)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def add_header(slide, title: str, section: str | None = None):
    add_textbox(slide, TITLE_X, TITLE_Y, TITLE_W, Inches(0.45), title, size=22, bold=True, color=NAVY)
    if section:
        add_textbox(slide, Inches(10.5), TITLE_Y, Inches(2.2), Inches(0.35), section, size=9, color=MID, align=PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TITLE_X, LINE_Y, Inches(12.25), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, text: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), FOOTER_Y, Inches(12.2), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(237, 243, 252)
    shape.line.color.rgb = RGBColor(210, 223, 245)
    add_textbox(slide, Inches(0.72), FOOTER_Y + Inches(0.03), Inches(11.85), Inches(0.31), text, size=12, bold=True, color=NAVY)


def blank_slide(prs, title: str, footer: str, section: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, title, section)
    add_footer(slide, footer)
    return slide


def add_image_fit(slide, image: Path, x, y, w, h):
    from PIL import Image

    with Image.open(image) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        final_w = w
        final_h = w / img_ratio
    else:
        final_h = h
        final_w = h * img_ratio
    left = x + (w - final_w) / 2
    top = y + (h - final_h) / 2
    return slide.shapes.add_picture(str(image), left, top, width=final_w, height=final_h)


def add_image_cover(slide, image: Path, x, y, w, h):
    # python-pptx has no crop-to-cover helper; fill width and accept vertical margin if needed.
    return add_image_fit(slide, image, x, y, w, h)


def add_card(slide, x, y, w, h, title: str, body: str, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 253)
    shape.line.color.rgb = RGBColor(218, 226, 238)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.07), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_textbox(slide, x + Inches(0.22), y + Inches(0.12), w - Inches(0.35), Inches(0.3), title, size=14, bold=True, color=NAVY)
    add_textbox(slide, x + Inches(0.22), y + Inches(0.52), w - Inches(0.35), h - Inches(0.62), body, size=11, color=DARK)


def add_flow_box(slide, x, y, w, h, text, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 249, 255)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.05), w - Inches(0.16), h - Inches(0.1), text, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=BLUE):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    line.line.end_arrowhead = True
    return line


def add_metric_table(slide, x, y, w, h, headers, rows):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 253) if r_idx % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = DARK
    return table


def read_csv_rows(path: Path):
    with path.open("r", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_image_cover(slide, AI_DIR / "cover_hybrid_defect_detection.png", Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6.2), SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = 8
    overlay.line.fill.background()
    add_textbox(slide, Inches(0.62), Inches(1.02), Inches(5.7), Inches(1.3), "超表面光学前端辅助的\n工业缺陷检测", size=30, bold=True, color=NAVY)
    add_textbox(slide, Inches(0.68), Inches(2.62), Inches(5.25), Inches(0.55), "Metasurface optical frontend + lightweight electronic backend", size=14, color=DARK)
    add_textbox(slide, Inches(0.68), Inches(6.55), Inches(5.25), Inches(0.35), "Fabric 二分类验证  |  DAGM / SegDecNet 分割闭环", size=11, bold=True, color=NAVY)


def slide_ai_image(prs, title, img, footer, section):
    slide = blank_slide(prs, title, footer, section)
    add_image_fit(slide, AI_DIR / img, Inches(0.75), CONTENT_TOP, Inches(11.85), Inches(5.55))
    return slide


def slide_route(prs):
    slide = blank_slide(prs, "总体研究路线", "整个课题沿着“电子 teacher -> optical student -> physical mapping”逐步推进。", "总体方案")
    add_image_fit(slide, AI_DIR / "research_route.png", Inches(0.62), Inches(1.1), Inches(6.1), Inches(4.9))
    add_flow_box(slide, Inches(7.0), Inches(1.25), Inches(1.75), Inches(0.7), "Electronic\nTeacher", NAVY)
    add_flow_box(slide, Inches(9.0), Inches(1.25), Inches(1.75), Inches(0.7), "Optical\nStudent", TEAL)
    add_flow_box(slide, Inches(11.0), Inches(1.25), Inches(1.75), Inches(0.7), "Physical\nProbe", BLUE)
    add_arrow(slide, Inches(8.75), Inches(1.6), Inches(9.0), Inches(1.6))
    add_arrow(slide, Inches(10.75), Inches(1.6), Inches(11.0), Inches(1.6))
    add_card(slide, Inches(7.0), Inches(2.25), Inches(5.75), Inches(0.85), "1. 从强电子模型出发", "先复现或固定 teacher，获得可靠任务监督。", NAVY)
    add_card(slide, Inches(7.0), Inches(3.25), Inches(5.75), Inches(0.85), "2. 设计物理友好 student", "前端限制为 optical convolution bank，后端保留必要非线性。", TEAL)
    add_card(slide, Inches(7.0), Inches(4.25), Inches(5.75), Inches(0.85), "3. 导出 kernels 并映射 PSF", "signed kernels 经过正负拆分，进入超表面可实现性验证。", BLUE)


def slide_kd_method(prs):
    slide = blank_slide(prs, "方法论：知识蒸馏连接软件模型和物理前端", "蒸馏用于把强电子模型的任务能力迁移到物理友好的学生模型。", "总体方案")
    add_flow_box(slide, Inches(0.9), Inches(1.55), Inches(2.4), Inches(1.0), "Teacher\nHigh-capacity CNN / SegDecNet", NAVY)
    add_flow_box(slide, Inches(5.45), Inches(1.55), Inches(2.4), Inches(1.0), "Student\nOptical frontend + backend", TEAL)
    add_arrow(slide, Inches(3.35), Inches(2.05), Inches(5.4), Inches(2.05))
    add_textbox(slide, Inches(3.65), Inches(1.55), Inches(1.55), Inches(0.35), "KD losses", size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(8.45), Inches(1.18), Inches(3.85), Inches(1.05), "Task loss", "真实标签监督 classification / segmentation。", BLUE)
    add_card(slide, Inches(8.45), Inches(2.48), Inches(3.85), Inches(1.05), "Soft prediction KD", "对齐 teacher 的分类概率、分割热图或 logits。", TEAL)
    add_card(slide, Inches(8.45), Inches(3.78), Inches(3.85), Inches(1.05), "Feature / volume KD", "在 DAGM 中对齐共享卷积 volume 和 segmentation 输出。", RED)
    add_textbox(slide, Inches(1.05), Inches(3.25), Inches(6.6), Inches(1.35), "训练目标示意：\nL = L_task + λ_seg L_segKD + λ_vol L_volumeKD + λ_cls L_clsKD", size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def slide_kernel_to_psf(prs):
    slide = blank_slide(prs, "从 learned kernels 到超表面 PSF", "signed kernel 需要做 positive/negative split 才能映射到非负光强响应。", "总体方案")
    labels = ["Signed kernel K", "Positive branch\nmax(K, 0)", "Negative branch\nmax(-K, 0)", "Target PSF", "Metasurface\nphase / radius"]
    xs = [0.7, 3.0, 5.3, 7.6, 9.9]
    colors = [NAVY, TEAL, RED, BLUE, NAVY]
    for i, (label, x, color) in enumerate(zip(labels, xs, colors)):
        add_flow_box(slide, Inches(x), Inches(2.05), Inches(1.65), Inches(1.0), label, color)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.68), Inches(2.55), Inches(xs[i + 1] - 0.05), Inches(2.55), color=BLUE)
    add_textbox(slide, Inches(1.15), Inches(4.1), Inches(11.0), Inches(0.8), "核心公式：K = K_positive - K_negative", size=24, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5), "这个拆分逻辑在 Fabric 和 DAGM 两个案例中都保持一致。", size=16, color=DARK, align=PP_ALIGN.CENTER)


def slide_two_cases(prs):
    slide = blank_slide(prs, "两个案例在整体路线中的位置", "Fabric 是方法起点，DAGM 是更完整的工业缺陷分割主线。", "总体方案")
    add_card(slide, Inches(0.85), Inches(1.35), Inches(5.6), Inches(3.9), "案例一：Fabric patch 二分类", "任务：AITEX fabric patch OK/NG\n模型：teacher CNN -> R1 optical student\n重点：64x64 低分辨率输入、一层 optical kernels、阈值校准\n角色：小而清晰的方法验证 demo", TEAL)
    add_card(slide, Inches(6.85), Inches(1.35), Inches(5.6), Inches(3.9), "案例二：DAGM / SegDecNet 分割", "任务：DAGM Class7 缺陷分类 + mask 分割\n模型：SegDecNet teacher -> optical student\n重点：两阶段蒸馏、mask 结果、PSF target、metasurface probe\n角色：当前高性能主结果闭环", BLUE)
    add_arrow(slide, Inches(6.45), Inches(3.3), Inches(6.82), Inches(3.3), NAVY)


def slide_fabric_task(prs):
    slide = blank_slide(prs, "Fabric / AITEX：patch 二分类任务", "Fabric 提供了一个易于讲清楚的一层 optical student 起点。", "Fabric")
    add_image_fit(slide, FABRIC / "figures_main" / "results" / "ui.png", Inches(0.7), Inches(1.22), Inches(5.5), Inches(4.8))
    x0 = Inches(6.65)
    add_flow_box(slide, x0, Inches(1.55), Inches(1.75), Inches(0.75), "Full fabric\nimage", NAVY)
    add_flow_box(slide, x0 + Inches(2.05), Inches(1.55), Inches(1.75), Inches(0.75), "256x256\npatches", TEAL)
    add_flow_box(slide, x0 + Inches(4.1), Inches(1.55), Inches(1.75), Inches(0.75), "OK / NG\nclassification", BLUE)
    add_arrow(slide, x0 + Inches(1.78), Inches(1.92), x0 + Inches(2.05), Inches(1.92))
    add_arrow(slide, x0 + Inches(3.83), Inches(1.92), x0 + Inches(4.1), Inches(1.92))
    add_card(slide, Inches(6.65), Inches(2.8), Inches(5.85), Inches(1.05), "任务特点", "缺陷稀疏、patch 二分类、适合快速验证 optical frontend。", TEAL)
    add_card(slide, Inches(6.65), Inches(4.1), Inches(5.85), Inches(1.05), "为什么先做它", "链路短、解释清楚，能快速检验低分辨率 student 是否有效。", BLUE)


def slide_architecture_image(prs, title, image, footer, section):
    slide = blank_slide(prs, title, footer, section)
    add_image_fit(slide, image, Inches(0.8), CONTENT_TOP, Inches(11.7), Inches(5.55))


def slide_fabric_teacher_arch(prs):
    slide = blank_slide(prs, "Fabric teacher CNN 与 patch pipeline", "Teacher 是固定电子基线，student 只替换早期特征提取思路。", "Fabric")
    labels = ["Fabric image\n256x4096", "Patch split\n16 x 256x256", "CNN feature\nextractor", "FC classifier", "OK / NG\nscore"]
    xs = [0.65, 2.95, 5.25, 7.55, 9.85]
    colors = [NAVY, TEAL, TEAL, BLUE, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(2.0), Inches(1.8), Inches(0.95), label, color)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.83), Inches(2.48), Inches(xs[i + 1] - 0.06), Inches(2.48))
    rows = [["Teacher reference", "F1 ≈ 0.975"], ["Patch size", "256x256"], ["Student input", "64x64"], ["Student target", "binary defect score"]]
    add_metric_table(slide, Inches(2.0), Inches(4.0), Inches(9.3), Inches(1.55), ["Item", "Value"], rows)


def slide_fabric_student_arch(prs):
    slide = blank_slide(prs, "Fabric R1 student：低分辨率 optical frontend", "关键经验是把输入压到 64x64，而不是盲目堆学生网络。", "Fabric")
    labels = ["Input\n64x64", "Optical conv\n16 x 7x7", "ReLU + Pool\npooled=6", "FC backend\nhidden=256", "Binary\nscore"]
    xs = [0.8, 3.0, 5.35, 7.7, 10.05]
    colors = [NAVY, TEAL, TEAL, BLUE, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(2.2), Inches(1.75), Inches(0.95), label, color)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.78), Inches(2.68), Inches(xs[i + 1] - 0.05), Inches(2.68))
    add_card(slide, Inches(1.25), Inches(4.15), Inches(10.7), Inches(1.1), "锁定配置", "input_size=64, optical_kernels=16, kernel_size=7, pooled_size=6, hidden_dim=256, activation=ReLU", BLUE)


def slide_single_image(prs, title, image, footer, section):
    slide = blank_slide(prs, title, footer, section)
    add_image_fit(slide, image, Inches(0.85), CONTENT_TOP, Inches(11.65), Inches(5.55))


def slide_two_images(prs, title, image_left, image_right, footer, section):
    slide = blank_slide(prs, title, footer, section)
    add_image_fit(slide, image_left, Inches(0.75), CONTENT_TOP, Inches(5.65), Inches(5.45))
    add_image_fit(slide, image_right, Inches(6.85), CONTENT_TOP, Inches(5.65), Inches(5.45))


def slide_fabric_compute(prs):
    slide = blank_slide(prs, "Fabric 计算量和硬件意义", "Fabric 证明早期特征提取可被压缩到更轻的光电混合形式。", "Fabric")
    add_image_fit(slide, FABRIC / "figures_main" / "compute" / "fabric_compute_comparison_bar.png", Inches(0.65), Inches(1.2), Inches(6.0), Inches(4.95))
    rows = [["Teacher CNN", "26.08M", "733.77M"], ["R1 student", "0.149M", "7.37M"], ["Reduction", "175x", "99.5x"]]
    add_metric_table(slide, Inches(7.0), Inches(2.0), Inches(5.1), Inches(2.2), ["Model", "Params", "MACs"], rows)
    add_card(slide, Inches(7.0), Inches(4.65), Inches(5.1), Inches(0.95), "汇报口径", "这些是理论 MAC/参数量估计，用于说明硬件意义，不等价于实测速度。", RED)


def slide_dagm_task(prs):
    slide = blank_slide(prs, "DAGM / SegDecNet：工业纹理缺陷分割任务", "DAGM 更接近真实工业纹理异常定位和分割任务。", "DAGM")
    add_image_fit(slide, DAGM / "figures_main" / "qualitative_masks" / "mask_visualization_contact_sheet_12.jpg", Inches(0.72), Inches(1.18), Inches(7.0), Inches(4.8))
    add_card(slide, Inches(8.0), Inches(1.35), Inches(4.45), Inches(1.15), "任务形式", "输入灰度工业纹理图像，输出 image-level defect score 与 pixel-level mask。", BLUE)
    add_card(slide, Inches(8.0), Inches(2.85), Inches(4.45), Inches(1.15), "为什么更关键", "它不仅要判断有没有缺陷，还要定位缺陷区域。", TEAL)
    add_card(slide, Inches(8.0), Inches(4.35), Inches(4.45), Inches(1.15), "本章角色", "这是当前主结果，展示完整蒸馏、分割、PSF 和硬件意义闭环。", NAVY)


def slide_segdec_teacher(prs):
    slide = blank_slide(prs, "原始 SegDecNet teacher：共享卷积 volume + 分割/分类头", "SegDecNet teacher 提供分割与分类联合监督，是 optical student 蒸馏的核心来源。", "DAGM")
    add_flow_box(slide, Inches(0.85), Inches(2.25), Inches(1.55), Inches(0.8), "Input\nimage", NAVY)
    add_flow_box(slide, Inches(2.75), Inches(2.1), Inches(2.0), Inches(1.1), "Shared conv\nbackbone volume", TEAL)
    add_flow_box(slide, Inches(5.25), Inches(1.35), Inches(1.8), Inches(0.85), "Segmentation\nhead", BLUE)
    add_flow_box(slide, Inches(5.25), Inches(3.05), Inches(1.8), Inches(0.85), "Feature\nextractor", BLUE)
    add_flow_box(slide, Inches(7.65), Inches(3.05), Inches(1.8), Inches(0.85), "FC\nclassifier", RED)
    add_flow_box(slide, Inches(10.05), Inches(2.2), Inches(1.8), Inches(0.85), "Mask +\ndefect score", NAVY)
    add_arrow(slide, Inches(2.4), Inches(2.65), Inches(2.75), Inches(2.65))
    add_arrow(slide, Inches(4.75), Inches(2.45), Inches(5.25), Inches(1.8))
    add_arrow(slide, Inches(4.75), Inches(2.8), Inches(5.25), Inches(3.48))
    add_arrow(slide, Inches(7.05), Inches(3.48), Inches(7.65), Inches(3.48))
    add_arrow(slide, Inches(9.45), Inches(3.48), Inches(10.05), Inches(2.65))
    add_arrow(slide, Inches(7.05), Inches(1.8), Inches(10.05), Inches(2.45))
    add_card(slide, Inches(1.2), Inches(4.7), Inches(10.9), Inches(0.85), "蒸馏重点", "student 不是只训练 segonly，而是保留 SegDecNet 的分割、特征提取和分类骨架，同时把早期 volume 前端改成 optical convolution bank。", TEAL)


def slide_dagm_student(prs):
    slide = blank_slide(prs, "DAGM optical student：单层 optical bank + 轻量电子后端", "光学前端承担早期卷积，电子后端保留关键非线性决策能力。", "DAGM")
    labels = ["Input\n256x256", "Optical conv bank\n64 x 15x15", "FeatureNorm\n+ ReLU", "AvgPool\nstride=4", "Seg head\n+ volume concat", "Extractor\n+ FC"]
    xs = [0.55, 2.45, 4.75, 6.75, 8.6, 10.55]
    colors = [NAVY, TEAL, TEAL, BLUE, BLUE, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(2.15), Inches(1.55), Inches(0.95), label, color)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.58), Inches(2.62), Inches(xs[i + 1] - 0.04), Inches(2.62))
    add_card(slide, Inches(1.0), Inches(4.25), Inches(5.25), Inches(1.0), "光学部分", "单层 convolution bank，是后续 PSF / metasurface 映射的主要对象。", TEAL)
    add_card(slide, Inches(6.65), Inches(4.25), Inches(5.25), Inches(1.0), "电子部分", "seg head、extractor、FC 参数不多但功能关键，因此保留骨架并适度简化。", BLUE)


def slide_two_stage(prs):
    slide = blank_slide(prs, "两阶段蒸馏训练策略", "先对齐 optical/segmentation，再联合训练全 student，更适合当前任务。", "DAGM")
    add_card(slide, Inches(0.85), Inches(1.45), Inches(5.55), Inches(2.0), "Stage 1: frontend / segmentation warm-up", "重点强化 segmentation distillation、volume KD 和 mask soft target。\n目标是先让 optical bank 学到与 teacher volume 对齐的前端响应。", TEAL)
    add_card(slide, Inches(6.95), Inches(1.45), Inches(5.55), Inches(2.0), "Stage 2: joint optimization", "联合 classification、segmentation、volume distillation 和 task loss。\n目标是恢复完整 SegDecNet-style 决策能力。", BLUE)
    add_arrow(slide, Inches(6.4), Inches(2.45), Inches(6.93), Inches(2.45))
    rows = [["seg KD", "high"], ["volume KD", "high"], ["task cls loss", "joint stage"], ["threshold", "0.50 stable"]]
    add_metric_table(slide, Inches(3.1), Inches(4.2), Inches(7.1), Inches(1.55), ["Component", "Role"], rows)


def slide_dagm_exploration(prs):
    slide = blank_slide(prs, "架构探索与最终选择", "最终选择在性能、参数量和物理友好性之间折中。", "DAGM")
    rows = [
        ["input", "256x256", "保留足够分割细节"],
        ["optical bank", "64 kernels", "容量比 32 更稳"],
        ["kernel", "15x15", "纹理/边缘感受野更强"],
        ["downsample", "AvgPool stride=4", "降低电子后端计算"],
        ["training", "two-stage KD", "比单一 segonly 更符合最终目标"],
    ]
    add_metric_table(slide, Inches(0.95), Inches(1.55), Inches(11.45), Inches(3.5), ["Design axis", "Final choice", "Reason"], rows)
    add_textbox(slide, Inches(1.2), Inches(5.4), Inches(10.9), Inches(0.45), "最终主线：256 / optical64 / kernel15 / downsample4 / two-stage distillation", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_dagm_metrics(prs):
    slide = blank_slide(prs, "DAGM 最优 student 定量结果", "Full validation 上 IoU=0.9145、Dice=0.9349，结果很强。", "DAGM")
    add_image_fit(slide, DAGM / "figures_main" / "metrics" / "dagm_full_validation_bar.png", Inches(0.7), Inches(1.15), Inches(6.4), Inches(4.85))
    rows = [["AP", "1.000"], ["AUC", "1.000"], ["IoU", "0.9145"], ["Dice", "0.9349"], ["Precision", "0.9958"], ["Recall", "0.9176"]]
    add_metric_table(slide, Inches(7.65), Inches(1.55), Inches(4.2), Inches(3.8), ["Metric", "Value"], rows)


def slide_dagm_compute(prs):
    slide = blank_slide(prs, "DAGM 计算量节省与硬件意义", "电子后端 MACs 可理论上降到 teacher 的很小比例，硬件意义明确。", "DAGM")
    add_image_fit(slide, DAGM / "figures_main" / "compute" / "compute_comparison_bar.png", Inches(0.65), Inches(1.2), Inches(5.85), Inches(4.85))
    rows = [["Params reduction", "338x"], ["Digital student MACs", "21.8x fewer"], ["Hybrid backend MACs", "905x fewer"], ["vs 512 teacher", "3618x fewer"]]
    add_metric_table(slide, Inches(7.0), Inches(1.85), Inches(5.15), Inches(2.75), ["Comparison", "Ratio"], rows)
    add_card(slide, Inches(7.0), Inches(5.0), Inches(5.15), Inches(0.75), "注意", "这是理论电子 MAC reduction，不等价于最终实测速度。", RED)


def slide_summary_compare(prs):
    slide = blank_slide(prs, "两个案例对照总结", "Fabric 验证方法起点，DAGM 验证高性能分割闭环。", "总结")
    rows = [
        ["Fabric", "patch 二分类", "R1 F1=0.8571", "低分辨率 optical student"],
        ["DAGM", "分类 + 分割", "IoU=0.9145, Dice=0.9349", "蒸馏到 PSF probe 闭环"],
    ]
    add_metric_table(slide, Inches(0.9), Inches(1.8), Inches(11.55), Inches(2.0), ["Case", "Task", "Best result", "Role"], rows)
    add_card(slide, Inches(1.1), Inches(4.45), Inches(10.95), Inches(1.0), "递进关系", "从 Fabric 的小型二分类 demo，到 DAGM 的分割主结果，课题已经形成从算法蒸馏到物理映射入口的连续证据链。", BLUE)


def slide_contributions(prs):
    slide = blank_slide(prs, "当前阶段已完成贡献", "课题已经从算法实验推进到超表面映射入口。", "总结")
    add_card(slide, Inches(0.8), Inches(1.55), Inches(3.75), Inches(3.65), "1. 模型蒸馏", "完成 Fabric 和 DAGM 两条 teacher-student 路线，明确 optical frontend 的任务能力边界。", NAVY)
    add_card(slide, Inches(4.8), Inches(1.55), Inches(3.75), Inches(3.65), "2. 光学前端设计", "导出 learned kernels，完成 signed / positive / negative split，并形成 PSF target。", TEAL)
    add_card(slide, Inches(8.8), Inches(1.55), Inches(3.75), Inches(3.65), "3. 物理可行性验证", "代表 kernel 的 metasurface probe 已达到较高 PSF 相似度，说明路线初步可行。", BLUE)


def slide_next_steps(prs):
    slide = blank_slide(prs, "下一步工作", "下一阶段重点是把 simulated PSF 放回模型链路评估性能下降。", "展望")
    rows = [
        ["Hardware-aware retraining", "加入 PSF 误差、噪声和 calibration"],
        ["Simulated PSF student", "用模拟 PSF 替代 digital kernels 重新评估"],
        ["Physical parameter sweep", "优化波长、距离、ROI、迭代次数和结构约束"],
        ["Real optical experiment", "搭建 metasurface + sensor 验证闭环"],
    ]
    add_metric_table(slide, Inches(1.05), Inches(1.55), Inches(11.1), Inches(3.45), ["Direction", "Why it matters"], rows)
    add_image_fit(slide, AI_DIR / "hybrid_optical_system.png", Inches(3.8), Inches(5.15), Inches(5.75), Inches(1.35))


def build_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_ai_image(prs, "工业缺陷检测的核心矛盾", "background_problem.png", "工业质检需要同时满足高分辨率、低延迟和低功耗。", "背景")
    slide_ai_image(prs, "本课题目标：把早期特征提取前移到光学链路", "hybrid_optical_system.png", "目标是在尽量保持任务性能的同时，用超表面承担部分前端卷积。", "背景")
    slide_route(prs)
    slide_kd_method(prs)
    slide_kernel_to_psf(prs)
    slide_two_cases(prs)

    slide_fabric_task(prs)
    slide_fabric_teacher_arch(prs)
    slide_fabric_student_arch(prs)
    slide_single_image(prs, "Fabric 主结果：teacher 与 R1 student", FABRIC / "figures_main" / "results" / "fabric_teacher_student_summary_bar.png", "R1 student 在最佳阈值下达到 F1=0.8571，已经具备原型价值。", "Fabric")
    slide_single_image(prs, "Fabric 阈值敏感性", FABRIC / "figures_process" / "threshold_sweep" / "fabric_r1_threshold_story.png", "默认阈值会低估模型，部署前需要阈值校准。", "Fabric")
    slide_two_images(prs, "Fabric optical kernels 与正负拆分", FABRIC / "figures_main" / "kernels" / "kernel_grid_signed.png", FABRIC / "figures_main" / "kernels" / "kernel_grid_positive.png", "Student kernels 可以导出，并转为 positive / negative PSF 目标。", "Fabric")
    slide_fabric_compute(prs)

    slide_dagm_task(prs)
    slide_segdec_teacher(prs)
    slide_dagm_student(prs)
    slide_two_stage(prs)
    slide_dagm_exploration(prs)
    slide_dagm_metrics(prs)
    slide_single_image(prs, "DAGM mask 定性结果", DAGM / "figures_main" / "qualitative_masks" / "mask_visualization_contact_sheet_12.jpg", "预测热图基本落在真实缺陷区域，mask 偏紧但误检少。", "DAGM")
    slide_single_image(prs, "DAGM threshold sweep", DAGM / "figures_process" / "threshold_sweep" / "dagm_threshold_sweep.png", "DAGM student 默认阈值 0.5 已经稳定，不是靠手调阈值刷分。", "DAGM")
    slide_two_images(prs, "DAGM optical kernels", DAGM / "figures_main" / "kernels" / "kernel_grid_signed.png", DAGM / "figures_main" / "kernels" / "kernel_grid_negative.png", "Learned kernels 呈现纹理、边缘和方向性结构，并需要 positive/negative split。", "DAGM")
    slide_two_images(prs, "DAGM PSF target 与 backphase", DAGM / "figures_main" / "psf_targets" / "psf_target_center_crop.png", DAGM / "figures_process" / "metasurface_probe" / "psf_backphase_preview.png", "Learned kernels 已经可以稳定转成单波长 target PSF 和初始 backphase。", "DAGM")
    slide_two_images(prs, "Metasurface feasibility probe", DAGM / "figures_main" / "metasurface_probe" / "kernel00_positive_probe.png", DAGM / "figures_main" / "metasurface_probe" / "metasurface_probe_cosine_bar.png", "代表 kernel 的 PSF 拟合相似度达到 0.979-0.991，物理可实现性初步成立。", "DAGM")
    slide_dagm_compute(prs)
    slide_single_image(prs, "DAGM 小结：从蒸馏到物理 probe", AI_DIR / "research_route.png", "DAGM 已形成“蒸馏 -> kernel -> PSF target -> feasibility probe”的完整链路。", "DAGM")

    slide_summary_compare(prs)
    slide_contributions(prs)
    slide_next_steps(prs)
    slide_ai_image(prs, "未来展望：从分类/分割到定位检测", "future_yolo_chip_inspection.png", "后续可扩展到 YOLO、芯片、PCB、wafer 等定位检测任务。", "展望")

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build_deck()
