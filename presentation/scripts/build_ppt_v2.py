from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "presentation" / "build"
FORMULA_DIR = ROOT / "presentation" / "generated_figures" / "formulas"
AI_DIR = ROOT / "presentation" / "ai_images"
AI_V2_DIR = ROOT / "presentation" / "ai_images_v2"
FABRIC = ROOT / "fabric_defect_detection-main" / "paper_assets"
DAGM = ROOT / "mixed-segdec-net-comind2021-master" / "paper_assets"

PPTX_PATH = OUT_DIR / "metasurface_industrial_defect_detection_v2.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TITLE_X = Inches(0.55)
TITLE_Y = Inches(0.26)
TITLE_W = Inches(12.25)
LINE_Y = Inches(0.88)
FOOTER_Y = Inches(6.86)

NAVY = RGBColor(22, 58, 108)
BLUE = RGBColor(47, 107, 255)
TEAL = RGBColor(42, 157, 143)
RED = RGBColor(231, 111, 81)
AMBER = RGBColor(230, 167, 72)
DARK = RGBColor(36, 50, 74)
MID = RGBColor(95, 111, 132)
PALE = RGBColor(247, 249, 252)
LINE = RGBColor(213, 222, 236)
WHITE = RGBColor(255, 255, 255)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    FORMULA_DIR.mkdir(parents=True, exist_ok=True)


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_title(slide, title: str, section: str | None = None):
    add_textbox(slide, TITLE_X, TITLE_Y, TITLE_W, Inches(0.5), title, size=22, bold=True, color=NAVY)
    if section:
        add_textbox(slide, Inches(10.6), TITLE_Y, Inches(2.1), Inches(0.35), section, size=9, color=MID, align=PP_ALIGN.RIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, TITLE_X, LINE_Y, Inches(12.25), Inches(0.025))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()


def add_footer(slide, text: str):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), FOOTER_Y, Inches(12.25), Inches(0.38))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(237, 243, 252)
    box.line.color.rgb = RGBColor(210, 223, 245)
    add_textbox(slide, Inches(0.72), FOOTER_Y + Inches(0.03), Inches(11.9), Inches(0.31), text, size=12, bold=True, color=NAVY)


def blank(prs, title: str, footer: str, section: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_title(slide, title, section)
    add_footer(slide, footer)
    return slide


def add_card(slide, x, y, w, h, title: str, body: str, *, accent=BLUE, title_size=13, body_size=10.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE
    shape.line.color.rgb = LINE
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    add_textbox(slide, x + Inches(0.18), y + Inches(0.08), w - Inches(0.3), Inches(0.32), title, size=title_size, bold=True, color=NAVY)
    add_textbox(slide, x + Inches(0.18), y + Inches(0.45), w - Inches(0.32), h - Inches(0.52), body, size=body_size, color=DARK, valign=MSO_ANCHOR.TOP)
    return shape


def add_flow_box(slide, x, y, w, h, text: str, *, color=BLUE, fill=RGBColor(245, 249, 255), size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = color
    shape.line.width = Pt(1.4)
    add_textbox(slide, x + Inches(0.05), y + Inches(0.03), w - Inches(0.1), h - Inches(0.06), text, size=size, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, *, color=BLUE, width=2.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_image_fit(slide, path: Path, x, y, w, h):
    from PIL import Image

    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        fw = w
        fh = w / img_ratio
    else:
        fh = h
        fw = h * img_ratio
    return slide.shapes.add_picture(str(path), x + (w - fw) / 2, y + (h - fh) / 2, width=fw, height=fh)


def add_image_frame(slide, path: Path, x, y, w, h):
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    add_image_fit(slide, path, x + Inches(0.08), y + Inches(0.08), w - Inches(0.16), h - Inches(0.16))


def add_table(slide, x, y, w, h, headers, rows, *, font_size=8.5):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for r_i, row in enumerate(rows, 1):
        for c_i, value in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 253) if r_i % 2 else WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = DARK
    return table


def formula(name: str, expr: str, *, fontsize=24, width=7.0, height=0.8) -> Path:
    path = FORMULA_DIR / f"{name}.png"
    fig = plt.figure(figsize=(width, height), dpi=220)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, expr, ha="center", va="center", fontsize=fontsize, color="#163A6C")
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)
    return path


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_textbox(slide, Inches(1.0), Inches(1.55), Inches(11.2), Inches(0.9), "超表面光学前端辅助的工业缺陷检测方法研究", size=30, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.35), Inches(2.55), Inches(10.6), Inches(0.45), "Metasurface-based Optical Frontend for Industrial Defect Detection", size=16, color=DARK, align=PP_ALIGN.CENTER)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.25), Inches(3.35), Inches(8.85), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    add_textbox(slide, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.45), "李升欣", size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.2), Inches(4.82), Inches(4.9), Inches(0.4), "2026年5月", size=15, color=DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.2), Inches(6.55), Inches(10.9), Inches(0.35), "Fabric 二分类验证  |  DAGM / SegDecNet 缺陷分割闭环  |  Kernel-to-PSF 物理映射", size=11, bold=True, color=MID, align=PP_ALIGN.CENTER)


def background_slide(prs):
    slide = blank(prs, "背景：工业缺陷检测与超表面机器视觉前端的交叉点", "工业质检需要高速低功耗视觉前端，而超表面提供了把早期卷积前移到成像链路的硬件入口。", "背景")
    add_image_frame(slide, AI_V2_DIR / "v2_background_infographic.png", Inches(0.62), Inches(1.18), Inches(4.35), Inches(3.6))
    add_card(slide, Inches(5.2), Inches(1.15), Inches(3.45), Inches(1.25), "工业检测痛点", "高分辨率图像带来大量 early convolution MACs；缺陷往往稀疏、细小，边缘部署还受功耗和延迟约束。", accent=RED)
    add_card(slide, Inches(8.95), Inches(1.15), Inches(3.65), Inches(1.25), "传统电子模型", "CNN classifier、U-Net/SegNet/SegDecNet、YOLO 等模型通常依赖电子端卷积前端。", accent=BLUE)
    add_card(slide, Inches(5.2), Inches(2.75), Inches(3.45), Inches(1.25), "超表面机会", "通过 PSF / optical convolution bank 在成像阶段编码早期特征，减少后端电子计算压力。", accent=TEAL)
    add_card(slide, Inches(8.95), Inches(2.75), Inches(3.65), Inches(1.25), "本课题切入点", "不是做全光学 detector，而是做 metasurface optical frontend + lightweight electronic backend。", accent=NAVY)
    add_textbox(slide, Inches(0.75), Inches(5.25), Inches(11.75), Inches(0.75), "交叉问题：超表面光学前端是否能为工业缺陷检测承担早期卷积特征提取，并在保持任务性能的同时降低电子计算量？", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def traditional_pipeline_slide(prs):
    slide = blank(prs, "传统工业视觉模型：性能来自电子卷积前端，但计算压力也集中在这里", "传统 CNN/YOLO/Segmentation 模型都依赖电子端 early feature extraction。", "方法")
    labels = ["Industrial\nimage", "Electronic\nCNN frontend", "Feature\nmaps", "Task head", "Output"]
    xs = [0.8, 3.05, 5.6, 7.85, 10.1]
    for i, (x, label) in enumerate(zip(xs, labels)):
        add_flow_box(slide, Inches(x), Inches(1.65), Inches(1.65), Inches(0.85), label, color=[NAVY, RED, BLUE, BLUE, TEAL][i])
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.68), Inches(2.08), Inches(xs[i + 1] - 0.06), Inches(2.08))
    add_card(slide, Inches(0.9), Inches(3.2), Inches(3.5), Inches(1.25), "分类", "Binary CNN / ResNet style classifier\n输出：defect / normal", accent=BLUE)
    add_card(slide, Inches(4.9), Inches(3.2), Inches(3.5), Inches(1.25), "分割", "U-Net / SegNet / SegDecNet\n输出：pixel-level defect mask", accent=TEAL)
    add_card(slide, Inches(8.9), Inches(3.2), Inches(3.5), Inches(1.25), "检测定位", "YOLO / detector style model\n输出：defect bounding boxes", accent=AMBER)
    add_textbox(slide, Inches(1.05), Inches(5.3), Inches(11.25), Inches(0.55), "瓶颈：高分辨率输入上的早期卷积层通常占据大量 MACs，且必须在电子端逐帧计算。", size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)


def hybrid_architecture_slide(prs):
    slide = blank(prs, "光电融合架构：用超表面承担早期卷积，用电子后端完成非线性决策", "本课题的核心是把最早期、最昂贵的特征提取前移到成像链路。", "方法")
    labels = ["Sample\nillumination", "Metasurface\nPSF bank", "CMOS sensor\nfeature volume", "Calibration\n/ normalization", "Electronic\nbackend", "Defect score\n/ mask"]
    xs = [0.45, 2.35, 4.35, 6.35, 8.45, 10.55]
    colors = [NAVY, TEAL, TEAL, BLUE, BLUE, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(1.5), Inches(1.45), Inches(0.9), label, color=color, size=9.5)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.48), Inches(1.95), Inches(xs[i + 1] - 0.04), Inches(1.95))
    add_image_frame(slide, AI_V2_DIR / "v2_hybrid_optical_architecture.png", Inches(0.65), Inches(3.0), Inches(5.1), Inches(2.5))
    add_card(slide, Inches(6.1), Inches(3.05), Inches(3.0), Inches(1.05), "光学端做什么", "实现 CNN-like convolution / PSF feature encoding，输出多通道光学响应。", accent=TEAL)
    add_card(slide, Inches(9.45), Inches(3.05), Inches(3.0), Inches(1.05), "电子端做什么", "负责 calibration、非线性、分割/分类/检测 head，保持任务表达能力。", accent=BLUE)
    add_card(slide, Inches(6.1), Inches(4.55), Inches(6.35), Inches(1.0), "表述边界", "我们不声称“超表面实现完整神经网络”，而是构建 hybrid optical-electronic defect detection。", accent=NAVY)


def distillation_method_slide(prs):
    slide = blank(prs, "蒸馏训练：把强电子 teacher 的任务能力迁移到 optical student", "蒸馏是连接高性能电子模型和物理友好 student 的核心桥梁。", "方法")
    add_flow_box(slide, Inches(0.75), Inches(1.45), Inches(2.2), Inches(0.85), "Frozen teacher\nCNN / SegDecNet", color=NAVY)
    add_flow_box(slide, Inches(0.75), Inches(3.2), Inches(2.2), Inches(0.85), "Trainable student\noptical + electronic", color=TEAL)
    add_arrow(slide, Inches(3.0), Inches(1.88), Inches(5.0), Inches(1.88), color=NAVY)
    add_arrow(slide, Inches(3.0), Inches(3.63), Inches(5.0), Inches(3.63), color=TEAL)
    add_flow_box(slide, Inches(5.1), Inches(1.45), Inches(2.2), Inches(0.85), "Teacher outputs\nlogits / masks / volume", color=NAVY)
    add_flow_box(slide, Inches(5.1), Inches(3.2), Inches(2.2), Inches(0.85), "Student outputs\nscore / mask / volume", color=TEAL)
    add_arrow(slide, Inches(7.35), Inches(2.75), Inches(8.45), Inches(2.75), color=RED)
    add_flow_box(slide, Inches(8.55), Inches(2.2), Inches(2.35), Inches(1.1), "KD + task\nlosses", color=RED)
    add_image_fit(slide, formula("kd_loss_v2", r"$\mathcal{L}=\mathcal{L}_{task}+\lambda_{cls}\mathcal{L}_{KD}^{cls}+\lambda_{seg}\mathcal{L}_{KD}^{seg}+\lambda_{vol}\mathcal{L}_{KD}^{vol}$", fontsize=25, width=7.8), Inches(2.7), Inches(4.75), Inches(8.3), Inches(0.75))
    add_card(slide, Inches(11.15), Inches(1.25), Inches(1.45), Inches(4.7), "Loss 含义", "task：真实标签\ncls：分类软输出\nseg：soft mask\nvol：中间特征 volume", accent=RED, title_size=11, body_size=8.5)


def kernel_to_metasurface_slide(prs):
    slide = blank(prs, "从 learned kernel 到超表面 PSF：正负拆分是物理映射入口", "光强非负约束决定了 signed kernels 必须拆成 positive / negative 两路。", "方法")
    xs = [0.6, 2.6, 4.55, 6.55, 8.6, 10.6]
    labels = ["Learned\nkernel K", "Positive\nK+", "Negative\nK-", "Target\nPSF", "Phase / radius\noptimization", "Simulated\nPSF"]
    colors = [NAVY, TEAL, RED, BLUE, AMBER, NAVY]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(1.45), Inches(1.45), Inches(0.85), label, color=color, size=9.5)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.48), Inches(1.88), Inches(xs[i + 1] - 0.05), Inches(1.88))
    add_image_fit(slide, formula("kernel_split_v2", r"$K^{+}=\max(K,0),\quad K^{-}=\max(-K,0),\quad K=K^{+}-K^{-}$", fontsize=26, width=7.5), Inches(2.8), Inches(2.75), Inches(7.8), Inches(0.72))
    add_image_frame(slide, DAGM / "figures_main" / "psf_targets" / "psf_target_center_crop.png", Inches(0.8), Inches(4.0), Inches(3.6), Inches(1.95))
    add_image_frame(slide, DAGM / "figures_main" / "metasurface_probe" / "kernel00_positive_probe.png", Inches(4.85), Inches(3.75), Inches(4.1), Inches(2.45))
    add_card(slide, Inches(9.35), Inches(4.05), Inches(3.0), Inches(1.75), "谨慎口径", "当前是 first-pass physical feasibility probe：证明 PSF 形状初步可拟合，还不是最终器件制造结果。", accent=AMBER)


def two_cases_slide(prs):
    slide = blank(prs, "两个实验案例：从二分类 demo 到分割主结果", "Fabric 验证低分辨率 optical student；DAGM 验证完整缺陷分割闭环。", "方法")
    add_card(slide, Inches(0.75), Inches(1.2), Inches(5.8), Inches(4.6), "案例一：Fabric / AITEX 织物缺陷二分类", "工业场景：织物表面缺陷检测\n输入：fabric 长条图像切成 patch\n输出：patch-level normal / defective\n模型：Binary CNN teacher -> R1 optical student\n作用：验证一层 optical frontend 和低分辨率输入是否可行", accent=TEAL, body_size=11)
    add_card(slide, Inches(6.9), Inches(1.2), Inches(5.8), Inches(4.6), "案例二：DAGM Class7 纹理缺陷分割", "工业场景：纹理表面异常/划痕类缺陷\n输入：灰度工业纹理图像\n输出：image-level defect score + pixel-level mask\n模型：SegDecNet teacher -> optical student\n作用：验证高性能分割、PSF target 和 metasurface probe 闭环", accent=BLUE, body_size=11)


def methodology_overview_slide(prs):
    slide = blank(prs, "方法总览：从传统电子模型到光电融合 student", "这张图对应本课题的整体方法主图：传统模型、蒸馏、光学卷积和物理映射在同一条链路中。", "方法")
    add_image_frame(slide, AI_V2_DIR / "v2_methodology_graphical_abstract.png", Inches(0.6), Inches(1.12), Inches(7.25), Inches(4.75))
    add_card(slide, Inches(8.15), Inches(1.25), Inches(4.25), Inches(1.05), "传统模型", "CNN / YOLO / SegDecNet 在电子端完成特征提取和任务 head。", accent=BLUE)
    add_card(slide, Inches(8.15), Inches(2.65), Inches(4.25), Inches(1.05), "蒸馏迁移", "Teacher 输出、mask 和中间 volume 监督 optical student。", accent=RED)
    add_card(slide, Inches(8.15), Inches(4.05), Inches(4.25), Inches(1.05), "物理映射", "Learned kernels 经正负拆分后转为 PSF / metasurface optimization target。", accent=TEAL)


def fabric_intro_slide(prs):
    slide = blank(prs, "Fabric / AITEX：织物表面缺陷的 patch 二分类", "这个案例检测的是织物 patch 是否存在缺陷，不输出缺陷具体位置。", "Fabric")
    add_image_frame(slide, AI_V2_DIR / "v2_fabric_case_diagram.png", Inches(0.65), Inches(1.18), Inches(4.45), Inches(3.95))
    add_flow_box(slide, Inches(5.45), Inches(1.3), Inches(1.65), Inches(0.75), "Long fabric\nimage", color=NAVY, size=9.5)
    add_flow_box(slide, Inches(7.35), Inches(1.3), Inches(1.65), Inches(0.75), "256x256\npatches", color=TEAL, size=9.5)
    add_flow_box(slide, Inches(9.25), Inches(1.3), Inches(1.65), Inches(0.75), "Binary\nclassifier", color=BLUE, size=9.5)
    add_flow_box(slide, Inches(11.15), Inches(1.3), Inches(1.25), Inches(0.75), "OK / NG", color=RED, size=9.5)
    add_arrow(slide, Inches(7.12), Inches(1.67), Inches(7.35), Inches(1.67))
    add_arrow(slide, Inches(9.02), Inches(1.67), Inches(9.25), Inches(1.67))
    add_arrow(slide, Inches(10.92), Inches(1.67), Inches(11.15), Inches(1.67))
    add_card(slide, Inches(5.45), Inches(2.55), Inches(3.2), Inches(1.25), "任务定义", "Patch-level binary classification：判断当前织物 patch 是否包含缺陷。", accent=BLUE)
    add_card(slide, Inches(9.05), Inches(2.55), Inches(3.35), Inches(1.25), "为什么适合作为第一例", "链路短、指标清楚，能快速检验 optical convolution bank 是否有判别能力。", accent=TEAL)
    add_card(slide, Inches(5.45), Inches(4.25), Inches(6.95), Inches(1.05), "注意", "Fabric 本章不讲 U-Net，不讲 mask 定位；它的角色是二分类 demo 和架构起点。", accent=RED)


def fabric_teacher_slide(prs):
    slide = blank(prs, "Fabric teacher：电子 CNN 二分类基线", "Teacher 提供强电子基线，student 学的是 patch defect / normal 判别能力。", "Fabric")
    labels = ["Patch\n256x256", "Conv blocks\nfeature extraction", "Flatten / FC", "Sigmoid\nscore", "Defect?\nYes / No"]
    xs = [0.7, 2.8, 5.15, 7.25, 9.35]
    colors = [NAVY, BLUE, BLUE, TEAL, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(1.55), Inches(1.75), Inches(0.88), label, color=color)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.78), Inches(1.99), Inches(xs[i + 1] - 0.05), Inches(1.99))
    add_table(slide, Inches(2.0), Inches(3.25), Inches(9.35), Inches(1.65), ["Item", "Value"], [["Teacher model", "Binary CNN classifier"], ["Reference F1", "≈ 0.975"], ["Student target", "distill / approximate binary decision"], ["Industrial meaning", "fast patch-level screening"]], font_size=10)
    add_card(slide, Inches(1.6), Inches(5.25), Inches(10.2), Inches(0.75), "后续替换点", "Student 不复制完整 teacher，而是用一层 optical convolution bank 承担前端响应，再接轻量电子后端。", accent=TEAL)


def fabric_student_slide(prs):
    slide = blank(prs, "Fabric R1 student：64×64 输入 + 单层 optical convolution bank", "最关键经验是降低输入分辨率后，浅层 optical student 变得可用。", "Fabric")
    labels = ["Input\n64x64", "OpticalConvBank\n16 kernels, 7x7", "ReLU + pooling\npooled=6", "FC backend\nhidden=256", "Sigmoid\nbinary score"]
    xs = [0.65, 2.65, 5.25, 7.55, 10.0]
    colors = [NAVY, TEAL, TEAL, BLUE, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(1.55), Inches(1.9), Inches(0.9), label, color=color, size=10)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.92), Inches(2.0), Inches(xs[i + 1] - 0.05), Inches(2.0))
    add_table(slide, Inches(1.05), Inches(3.25), Inches(5.35), Inches(2.05), ["Parameter", "Value"], [["input_size", "64"], ["optical_kernels", "16"], ["kernel_size", "7"], ["pooled_size", "6"], ["hidden_dim", "256"]], font_size=9.5)
    add_card(slide, Inches(6.9), Inches(3.35), Inches(5.1), Inches(1.9), "设计解释", "光学端只保留一层 convolution bank，贴近未来 PSF/超表面实现；电子端用小 FC backend 完成非线性判别。", accent=TEAL)


def fabric_results_slide(prs):
    slide = blank(prs, "Fabric 结果：R1 student 需要阈值校准", "校准后 F1=0.8571，说明模型已学到有效判别特征。", "Fabric")
    add_image_frame(slide, FABRIC / "figures_main" / "results" / "fabric_teacher_student_summary_bar.png", Inches(0.65), Inches(1.18), Inches(6.0), Inches(3.95))
    add_image_frame(slide, FABRIC / "figures_process" / "threshold_sweep" / "fabric_r1_threshold_story.png", Inches(6.95), Inches(1.18), Inches(5.65), Inches(3.95))
    add_table(slide, Inches(1.1), Inches(5.45), Inches(11.0), Inches(0.85), ["Metric", "Default threshold", "Best threshold"], [["F1", "0.2353", "0.8571"], ["Precision / Recall", "not main result", "0.75 / 1.00"]], font_size=9)


def fabric_kernels_slide(prs):
    slide = blank(prs, "Fabric learned kernels：从二分类 student 到 PSF 输入", "Student kernels 可以导出，并转为 positive / negative PSF 目标。", "Fabric")
    add_image_frame(slide, FABRIC / "figures_main" / "kernels" / "kernel_grid_signed.png", Inches(0.65), Inches(1.15), Inches(3.75), Inches(3.75))
    add_image_frame(slide, FABRIC / "figures_main" / "kernels" / "kernel_grid_positive.png", Inches(4.75), Inches(1.15), Inches(3.75), Inches(3.75))
    add_image_frame(slide, FABRIC / "figures_main" / "kernels" / "kernel_grid_negative.png", Inches(8.85), Inches(1.15), Inches(3.75), Inches(3.75))
    add_card(slide, Inches(0.9), Inches(5.3), Inches(11.4), Inches(0.8), "观察", "Signed kernels 同时包含正负响应；由于光强非负，后续超表面映射需要 positive / negative 双路 PSF。", accent=TEAL)


def fabric_compute_slide(prs):
    slide = blank(prs, "Fabric 计算量：小型 student 显著压缩电子计算", "Fabric 是小而清晰的 demo，验证了压缩 early feature extraction 的方法价值。", "Fabric")
    add_image_frame(slide, FABRIC / "figures_main" / "compute" / "fabric_compute_comparison_bar.png", Inches(0.65), Inches(1.15), Inches(5.85), Inches(4.25))
    add_table(slide, Inches(6.9), Inches(1.55), Inches(5.15), Inches(2.4), ["Model", "Params", "MACs"], [["Teacher CNN", "26.08M", "733.77M"], ["R1 student", "0.149M", "7.37M"], ["Reduction", "≈175×", "≈99.5×"]], font_size=10)
    add_card(slide, Inches(6.9), Inches(4.35), Inches(5.15), Inches(0.95), "汇报口径", "这是理论参数/MACs 对比，用于说明硬件意义；最终速度还取决于曝光、读出、ADC 和电子实现。", accent=RED)


def fabric_summary_slide(prs):
    slide = blank(prs, "Fabric 小结：方法起点，而不是最终主结果", "Fabric 证明低分辨率 optical student 这条路能跑通，并为 DAGM 主结果提供经验。", "Fabric")
    add_card(slide, Inches(0.9), Inches(1.35), Inches(3.55), Inches(3.7), "1. 任务层面", "织物 patch 二分类，输出 defect / normal，不做缺陷位置分割。", accent=BLUE)
    add_card(slide, Inches(4.9), Inches(1.35), Inches(3.55), Inches(3.7), "2. 模型层面", "R1 student 使用 64x64 输入和单层 optical convolution bank，最佳阈值下 F1=0.8571。", accent=TEAL)
    add_card(slide, Inches(8.9), Inches(1.35), Inches(3.55), Inches(3.7), "3. 硬件层面", "Kernels 已导出并完成正负拆分，能够进入 PSF/metasurface 映射链路。", accent=NAVY)


def dagm_intro_slide(prs):
    slide = blank(prs, "DAGM Class7：工业纹理缺陷的分类 + 像素级分割", "这个案例不仅判断有没有缺陷，还输出缺陷区域 mask，是当前主结果。", "DAGM")
    add_image_frame(slide, AI_V2_DIR / "v2_dagm_case_diagram.png", Inches(0.65), Inches(1.15), Inches(6.65), Inches(4.35))
    add_card(slide, Inches(7.65), Inches(1.25), Inches(4.75), Inches(1.0), "工业场景", "DAGM 是工业纹理表面缺陷 benchmark；Class7 可视作纹理异常/划痕类检测场景。", accent=BLUE)
    add_card(slide, Inches(7.65), Inches(2.6), Inches(4.75), Inches(1.0), "输入输出", "输入灰度纹理图像；输出 image-level defect score 和 pixel-level defect mask。", accent=TEAL)
    add_card(slide, Inches(7.65), Inches(3.95), Inches(4.75), Inches(1.0), "为什么是主结果", "它比 Fabric 更复杂，能验证蒸馏、分割、kernel、PSF 和硬件意义的完整闭环。", accent=NAVY)


def segdec_teacher_slide(prs):
    slide = blank(prs, "SegDecNet teacher：共享卷积 volume 同时服务分割与分类", "Teacher 不只是 segmentation，它提供 mask 和 defect score 的联合监督。", "DAGM")
    add_flow_box(slide, Inches(0.65), Inches(2.2), Inches(1.45), Inches(0.8), "Input\nimage", color=NAVY)
    add_flow_box(slide, Inches(2.55), Inches(2.0), Inches(2.1), Inches(1.2), "Shared conv\nbackbone volume", color=TEAL)
    add_flow_box(slide, Inches(5.25), Inches(1.2), Inches(1.8), Inches(0.8), "Segmentation\nhead", color=BLUE)
    add_flow_box(slide, Inches(5.25), Inches(3.05), Inches(1.8), Inches(0.8), "Feature\nextractor", color=BLUE)
    add_flow_box(slide, Inches(7.65), Inches(3.05), Inches(1.7), Inches(0.8), "FC\nclassifier", color=RED)
    add_flow_box(slide, Inches(10.0), Inches(1.75), Inches(1.85), Inches(0.85), "Mask +\ndefect score", color=NAVY)
    add_arrow(slide, Inches(2.12), Inches(2.6), Inches(2.55), Inches(2.6))
    add_arrow(slide, Inches(4.68), Inches(2.35), Inches(5.25), Inches(1.6))
    add_arrow(slide, Inches(4.68), Inches(2.75), Inches(5.25), Inches(3.45))
    add_arrow(slide, Inches(7.08), Inches(1.6), Inches(10.0), Inches(2.0))
    add_arrow(slide, Inches(7.08), Inches(3.45), Inches(7.65), Inches(3.45))
    add_arrow(slide, Inches(9.35), Inches(3.45), Inches(10.0), Inches(2.35))
    add_card(slide, Inches(1.1), Inches(4.75), Inches(10.9), Inches(0.9), "蒸馏目标", "Student 需要同时继承 teacher 的分类能力、mask 监督和 shared volume 表征，而不是只训练 segonly。", accent=RED)


def dagm_student_slide(prs):
    slide = blank(prs, "DAGM optical student：超表面前端 + SegDecNet-style 电子后端", "光学前端承担早期卷积，电子后端保留分割和分类的关键骨架。", "DAGM")
    labels = ["Input\n256x256 gray", "OpticalConvBank\n64 kernels, 15x15", "FeatureNorm\n+ ReLU", "AvgPool\nkernel=4,stride=4", "Seg head\n1x1 conv", "Concat\nvolume + mask", "Extractor\n+ FC"]
    xs = [0.35, 1.95, 4.0, 5.75, 7.6, 9.15, 10.85]
    widths = [1.25, 1.65, 1.35, 1.45, 1.2, 1.35, 1.25]
    colors = [NAVY, TEAL, TEAL, BLUE, BLUE, AMBER, RED]
    for i, (x, width, label, color) in enumerate(zip(xs, widths, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(1.45), Inches(width), Inches(0.9), label, color=color, size=8.7)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + width + 0.03), Inches(1.9), Inches(xs[i + 1] - 0.04), Inches(1.9))
    add_card(slide, Inches(0.85), Inches(3.1), Inches(5.65), Inches(1.3), "未来光学实现部分", "OpticalConvBank 是主要映射对象：64 个 15x15 卷积核后续拆分为 positive / negative PSF 目标。", accent=TEAL)
    add_card(slide, Inches(6.9), Inches(3.1), Inches(5.55), Inches(1.3), "电子实现部分", "FeatureNorm、seg head、extractor、FC 参数不多但功能关键，负责非线性、mask 和分类决策。", accent=BLUE)
    add_table(slide, Inches(1.0), Inches(5.0), Inches(11.25), Inches(0.75), ["Optical bank", "Seg output", "Classifier output", "Main result"], [["64 x 15x15", "pixel-level mask", "image-level defect score", "IoU=0.9145 / Dice=0.9349"]], font_size=9)


def dagm_training_slide(prs):
    slide = blank(prs, "两阶段蒸馏：先对齐前端/分割，再联合优化完整任务", "两阶段训练比单一 segonly 更符合最终光电混合目标。", "DAGM")
    add_card(slide, Inches(0.75), Inches(1.25), Inches(5.65), Inches(1.55), "Stage 1: optical / segmentation warm-up", "强化 volume KD 和 segmentation KD，让 optical bank 先学到 teacher shared volume 和 soft mask。", accent=TEAL, body_size=11)
    add_card(slide, Inches(6.95), Inches(1.25), Inches(5.65), Inches(1.55), "Stage 2: joint optimization", "联合 task loss、classification KD、segmentation KD 和 volume KD，恢复完整 SegDecNet-style 决策能力。", accent=BLUE, body_size=11)
    add_arrow(slide, Inches(6.42), Inches(2.0), Inches(6.93), Inches(2.0), color=NAVY)
    add_image_fit(slide, formula("dagm_stage_loss_v2", r"$\mathcal{L}_{stage2}=\mathcal{L}_{task}+\lambda_{cls}\mathcal{L}_{KD}^{cls}+\lambda_{seg}\mathcal{L}_{KD}^{seg}+\lambda_{vol}\mathcal{L}_{KD}^{vol}$", fontsize=24, width=8.5), Inches(2.4), Inches(3.35), Inches(8.75), Inches(0.78))
    add_table(slide, Inches(2.0), Inches(4.65), Inches(9.3), Inches(1.05), ["Loss", "Purpose"], [["classification KD", "对齐 defect score / logits"], ["segmentation KD", "对齐 soft mask"], ["volume KD", "对齐 shared feature volume"]], font_size=9)


def dagm_exploration_slide(prs):
    slide = blank(prs, "架构探索：围绕光学容量、分辨率和蒸馏权重做折中", "最终选择在性能、参数量和物理友好性之间折中。", "DAGM")
    rows = [
        ["Optical bank", "32 -> 64", "64 更稳，仍保持单层光学前端"],
        ["Kernel size", "15x15", "更适合纹理/边缘响应"],
        ["Downsample", "AvgPool stride=4", "控制电子后端计算量"],
        ["Distillation", "seg_kd + volume_kd 加强", "让 mask 与前端 volume 更接近 teacher"],
        ["Training", "two-stage", "先 warm-up 再联合优化"],
    ]
    add_table(slide, Inches(0.75), Inches(1.35), Inches(11.85), Inches(3.75), ["Axis", "Choice", "Reason"], rows, font_size=8.8)
    add_card(slide, Inches(1.2), Inches(5.45), Inches(10.85), Inches(0.75), "最终主线", "256x256 input / optical64 / kernel15 / downsample4 / two-stage KD", accent=NAVY)


def dagm_metrics_slide(prs):
    slide = blank(prs, "DAGM 定量结果：分类满分，分割 IoU/Dice 较高", "Full validation 上 IoU=0.9145、Dice=0.9349，且 threshold=0.5 稳定。", "DAGM")
    add_image_frame(slide, DAGM / "figures_main" / "metrics" / "dagm_full_validation_bar.png", Inches(0.65), Inches(1.15), Inches(6.15), Inches(4.25))
    rows = [["AP", "1.000"], ["AUC", "1.000"], ["IoU", "0.9145"], ["Dice", "0.9349"], ["Precision", "0.9958"], ["Recall", "0.9176"]]
    add_table(slide, Inches(7.25), Inches(1.45), Inches(4.55), Inches(3.2), ["Metric", "Value"], rows, font_size=10)
    add_card(slide, Inches(7.25), Inches(5.0), Inches(4.55), Inches(0.75), "解释", "Precision 高于 Recall，说明 mask 偏保守、偏紧，但误检很少。", accent=TEAL)


def dagm_masks_slide(prs):
    slide = blank(prs, "DAGM mask 定性结果：能定位缺陷区域", "预测热图基本落在真实缺陷区域，mask 偏紧但误检少。", "DAGM")
    add_image_frame(slide, DAGM / "figures_main" / "qualitative_masks" / "mask_visualization_contact_sheet_12.jpg", Inches(0.65), Inches(1.15), Inches(7.2), Inches(4.7))
    add_card(slide, Inches(8.2), Inches(1.35), Inches(4.1), Inches(1.1), "看什么", "每组样例包含 input / GT mask / predicted heatmap / overlay。", accent=BLUE)
    add_card(slide, Inches(8.2), Inches(2.85), Inches(4.1), Inches(1.1), "观察一", "预测区域基本贴近 GT 缺陷位置，说明不是随机激活。", accent=TEAL)
    add_card(slide, Inches(8.2), Inches(4.35), Inches(4.1), Inches(1.1), "观察二", "mask 相对偏紧，与高 Precision、稍低 Recall 的定量结果一致。", accent=RED)


def dagm_threshold_slide(prs):
    slide = blank(prs, "DAGM threshold sweep：默认阈值 0.5 已稳定", "与 Fabric 不同，DAGM 主结果不是靠手调阈值刷出来的。", "DAGM")
    add_image_frame(slide, DAGM / "figures_process" / "threshold_sweep" / "dagm_threshold_sweep.png", Inches(0.8), Inches(1.15), Inches(7.5), Inches(4.7))
    add_card(slide, Inches(8.65), Inches(1.45), Inches(3.7), Inches(1.2), "Fabric", "二分类 student 阈值敏感，最佳阈值约 0.9。", accent=RED)
    add_card(slide, Inches(8.65), Inches(3.05), Inches(3.7), Inches(1.2), "DAGM", "默认 threshold=0.5 已经是稳定工作点。", accent=TEAL)
    add_card(slide, Inches(8.65), Inches(4.65), Inches(3.7), Inches(0.85), "汇报意义", "这让 DAGM 结果更适合作为主结果展示。", accent=NAVY)


def dagm_kernels_slide(prs):
    slide = blank(prs, "DAGM learned kernels：纹理、边缘和方向响应不是随机噪声", "Learned kernels 呈现可解释结构，并需要 positive/negative split。", "DAGM")
    add_image_frame(slide, DAGM / "figures_main" / "kernels" / "kernel_grid_signed.png", Inches(0.6), Inches(1.15), Inches(4.05), Inches(4.05))
    add_image_frame(slide, DAGM / "figures_main" / "kernels" / "kernel_grid_positive.png", Inches(4.85), Inches(1.15), Inches(3.75), Inches(4.05))
    add_image_frame(slide, DAGM / "figures_main" / "kernels" / "kernel_grid_negative.png", Inches(8.85), Inches(1.15), Inches(3.75), Inches(4.05))
    add_card(slide, Inches(0.95), Inches(5.55), Inches(11.4), Inches(0.65), "结论", "这些 kernels 是后续 PSF target 和 metasurface optimization 的直接输入。", accent=TEAL)


def dagm_psf_slide(prs):
    slide = blank(prs, "DAGM PSF target：把 learned kernels 转成可优化的光学目标", "Learned kernels 已经可以稳定转成单波长 target PSF 和初始 backphase。", "DAGM")
    add_image_fit(slide, formula("psf_mapping_v2", r"$K \rightarrow (K^{+},K^{-}) \rightarrow I_{target}(x,y) \rightarrow \phi(x,y) \rightarrow I_{sim}(x,y)$", fontsize=24, width=8.3), Inches(2.4), Inches(1.2), Inches(8.5), Inches(0.65))
    add_image_frame(slide, DAGM / "figures_main" / "psf_targets" / "psf_target_center_crop.png", Inches(0.75), Inches(2.25), Inches(5.6), Inches(3.65))
    add_image_frame(slide, DAGM / "figures_process" / "metasurface_probe" / "psf_backphase_preview.png", Inches(6.75), Inches(2.25), Inches(5.6), Inches(3.65))


def dagm_probe_slide(prs):
    slide = blank(prs, "Metasurface feasibility probe：代表 kernel 的 PSF 形状可拟合", "当前是 first-pass physical feasibility probe，不是最终器件制造结果。", "DAGM")
    add_image_frame(slide, DAGM / "figures_main" / "metasurface_probe" / "kernel00_positive_probe.png", Inches(0.65), Inches(1.15), Inches(6.1), Inches(4.8))
    add_image_frame(slide, DAGM / "figures_main" / "metasurface_probe" / "metasurface_probe_cosine_bar.png", Inches(7.15), Inches(1.15), Inches(5.25), Inches(3.55))
    add_card(slide, Inches(7.15), Inches(5.05), Inches(5.25), Inches(0.85), "结果", "代表 kernel 的 cosine similarity 在 0.979-0.991，说明目标 PSF 主结构初步可实现。", accent=TEAL)


def dagm_compute_slide(prs):
    slide = blank(prs, "DAGM 计算量节省：光学前端的硬件意义", "理论上 hybrid electronic backend 相对 teacher 的电子计算可减少到百倍量级以上。", "DAGM")
    add_image_frame(slide, DAGM / "figures_main" / "compute" / "compute_comparison_bar.png", Inches(0.65), Inches(1.15), Inches(5.85), Inches(4.3))
    rows = [["Teacher @256", "21.08G MACs"], ["Hybrid backend", "23.3M electronic MACs"], ["vs teacher@256", "905x fewer electronic MACs"], ["vs teacher@512", "3618x fewer electronic MACs"]]
    add_table(slide, Inches(6.95), Inches(1.45), Inches(5.2), Inches(2.65), ["Comparison", "Value"], rows, font_size=9.5)
    add_card(slide, Inches(6.95), Inches(4.6), Inches(5.2), Inches(0.9), "注意", "这是理论电子 MAC reduction；实际速度还取决于曝光、读出、ADC、光通量和电子硬件。", accent=RED)


def dagm_summary_slide(prs):
    slide = blank(prs, "DAGM 小结：已经形成从蒸馏到物理 probe 的闭环", "DAGM 路线承担当前高性能主结果角色。", "DAGM")
    labels = ["SegDecNet\nteacher", "Optical\nstudent", "High-IoU\nmask", "Kernel\nsplit", "PSF\ntarget", "Metasurface\nprobe", "Compute\nreduction"]
    xs = [0.35, 2.1, 3.85, 5.6, 7.35, 9.1, 10.85]
    colors = [NAVY, TEAL, BLUE, BLUE, AMBER, TEAL, RED]
    for i, (x, label, color) in enumerate(zip(xs, labels, colors)):
        add_flow_box(slide, Inches(x), Inches(2.0), Inches(1.25), Inches(0.85), label, color=color, size=8.8)
        if i < len(labels) - 1:
            add_arrow(slide, Inches(x + 1.28), Inches(2.42), Inches(xs[i + 1] - 0.05), Inches(2.42))
    add_table(slide, Inches(1.35), Inches(4.15), Inches(10.6), Inches(1.3), ["Result", "Value"], [["Segmentation", "IoU=0.9145, Dice=0.9349"], ["Metasurface probe", "cosine similarity 0.979-0.991"], ["Electronic MAC reduction", "905x fewer vs teacher@256"]], font_size=10)


def comparison_summary_slide(prs):
    slide = blank(prs, "两个案例对照：从可行 demo 到主结果闭环", "Fabric 验证方法起点，DAGM 验证高性能分割闭环。", "总结")
    rows = [
        ["Fabric / AITEX", "织物 patch 缺陷", "二分类", "F1=0.8571", "低分辨率 optical student"],
        ["DAGM Class7", "工业纹理异常", "分类 + 像素分割", "IoU=0.9145 / Dice=0.9349", "蒸馏 -> PSF -> probe 闭环"],
    ]
    add_table(slide, Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.85), ["Case", "Defect scene", "Output", "Main metric", "Role"], rows, font_size=8.5)
    add_card(slide, Inches(1.0), Inches(4.25), Inches(3.55), Inches(1.05), "算法贡献", "面向工业缺陷检测设计 optical frontend student。", accent=BLUE)
    add_card(slide, Inches(4.9), Inches(4.25), Inches(3.55), Inches(1.05), "系统贡献", "明确光学卷积 + 电子后端的混合架构边界。", accent=TEAL)
    add_card(slide, Inches(8.8), Inches(4.25), Inches(3.55), Inches(1.05), "物理贡献", "完成 kernel 到 PSF / metasurface probe 的入口验证。", accent=NAVY)


def future_slide(prs):
    slide = blank(prs, "未来工作：从 digital student 走向真实光电闭环", "下一阶段重点是把 simulated PSF 放回模型链路评估性能下降。", "展望")
    add_image_frame(slide, AI_V2_DIR / "v2_future_system_roadmap.png", Inches(0.65), Inches(1.18), Inches(4.65), Inches(3.25))
    rows = [
        ["Simulated PSF student", "用 simulated PSF 替代 learned kernels 重新评估性能"],
        ["Hardware-aware retraining", "加入噪声、偏移、制造误差和 calibration"],
        ["Physical sweep", "优化 wavelength、distance、ROI、iterations 和结构约束"],
        ["Detection extension", "扩展到 YOLO / PCB / wafer / chip defect localization"],
    ]
    add_table(slide, Inches(5.65), Inches(1.35), Inches(6.65), Inches(3.0), ["Direction", "Content"], rows, font_size=8.8)
    add_card(slide, Inches(0.95), Inches(5.25), Inches(11.4), Inches(0.75), "最终目标", "高速、低功耗、可部署的工业质检光电混合系统。", accent=NAVY)


def build() -> None:
    ensure_dirs()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover(prs)
    background_slide(prs)
    traditional_pipeline_slide(prs)
    methodology_overview_slide(prs)
    hybrid_architecture_slide(prs)
    distillation_method_slide(prs)
    kernel_to_metasurface_slide(prs)
    two_cases_slide(prs)

    fabric_intro_slide(prs)
    fabric_teacher_slide(prs)
    fabric_student_slide(prs)
    fabric_results_slide(prs)
    fabric_kernels_slide(prs)
    fabric_compute_slide(prs)
    fabric_summary_slide(prs)

    dagm_intro_slide(prs)
    segdec_teacher_slide(prs)
    dagm_student_slide(prs)
    dagm_training_slide(prs)
    dagm_exploration_slide(prs)
    dagm_metrics_slide(prs)
    dagm_masks_slide(prs)
    dagm_threshold_slide(prs)
    dagm_kernels_slide(prs)
    dagm_psf_slide(prs)
    dagm_probe_slide(prs)
    dagm_compute_slide(prs)
    dagm_summary_slide(prs)

    comparison_summary_slide(prs)
    future_slide(prs)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build()
